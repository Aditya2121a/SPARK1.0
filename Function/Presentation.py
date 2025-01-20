import pyttsx3
from pptx import Presentation
import time
def speak(text, pitch=50, rate=150, volume=0.9):
    engine = pyttsx3.init()
    engine.setProperty('rate', rate) # Speed of speech
    engine.setProperty('volume',volume) # Volume level (0.0 to 1.0)
    engine.setProperty('pitch', pitch) # Pitch level (0 to 100)

       # Real-time speaking
    engine.say(text)
    engine.runAndWait()

 # Create a presentation
presentation = Presentation()

slide_layout = presentation.slide_layouts[8]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Hello, Sir!"
subtitle.text = "This is an advanced speak function in your PowerPoint presentation. "
speak("Hello, Sir! This is an advanced speak function in your PowerPoint presentation.", pitch=60, rate=160, volume=0.8)
presentation.save("advanced_presentation.pptx")
time.sleep(2)
